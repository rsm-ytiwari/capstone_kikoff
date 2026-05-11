import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from templates import add_content_slide, _add_textbox, NAVY, GOLD, CHARCOAL
from pptx.util import Inches

_HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PROJECT_ROOT = os.path.abspath(os.path.join(_HERE, '..', '..'))
EDA_FIGS = os.path.join(PROJECT_ROOT, 'outputs', 'P1_02_eda', 'figures')

LEFT_CHART  = os.path.join(EDA_FIGS, 'monthly_seasonality.png')
RIGHT_CHART = os.path.join(EDA_FIGS, 'conversions_and_revenue.png')

CALLOUTS = [
    'Conversions peak in tax season (March–April); trough Nov–Dec — consistent with credit product demand cycle',
    'Spend-to-conversion lag = 0 days (r = 0.45) — media effects materialize quickly; adstock decay is short',
]


def build(prs):
    slide = add_content_slide(
        prs,
        'Results & Findings to Date — Seasonality and Conversion Patterns',
        page_num=10
    )

    if os.path.exists(LEFT_CHART):
        slide.shapes.add_picture(LEFT_CHART,
                                  Inches(0.40), Inches(1.58),
                                  Inches(5.90), Inches(4.70))

    if os.path.exists(RIGHT_CHART):
        slide.shapes.add_picture(RIGHT_CHART,
                                  Inches(6.60), Inches(1.58),
                                  Inches(6.10), Inches(4.70))

    for i, callout in enumerate(CALLOUTS):
        _add_textbox(slide, '• ' + callout,
                     left=0.50, top=6.38 + i * 0.52, width=12.20, height=0.50,
                     font_size=12, color=CHARCOAL)
