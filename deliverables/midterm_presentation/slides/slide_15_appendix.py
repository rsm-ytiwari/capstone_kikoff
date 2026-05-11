import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from templates import add_content_slide, _add_rect, _add_textbox, NAVY, GOLD, CHARCOAL, WHITE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

_HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PROJECT_ROOT = os.path.abspath(os.path.join(_HERE, '..', '..'))
PYMC_FIGS   = os.path.join(PROJECT_ROOT, 'outputs', 'P2_01_pymc', 'figures')
EDA_FIGS    = os.path.join(PROJECT_ROOT, 'outputs', 'P1_02_eda', 'figures')


# ── A1: Posterior Predictive Check ────────────────────────────────────────────
def _a1(prs):
    slide = add_content_slide(prs, 'Appendix A1 — Model Fit: Posterior Predictive Check',
                              page_num='A1')
    chart = os.path.join(PYMC_FIGS, 'posterior_predictive_check.png')
    if os.path.exists(chart):
        slide.shapes.add_picture(chart, Inches(0.40), Inches(1.58),
                                 Inches(12.40), Inches(5.20))
    _add_textbox(slide,
                 'Blue band = 94% HDI of posterior predictive distribution  ·  '
                 'Orange line = observed conversions  ·  Good overlap = adequate model fit',
                 left=0.50, top=6.85, width=12.20, height=0.40,
                 font_size=11, color=CHARCOAL)


# ── A2: Trace Plots ────────────────────────────────────────────────────────────
def _a2(prs):
    slide = add_content_slide(prs, 'Appendix A2 — MCMC Convergence: Trace Plots',
                              page_num='A2')
    chart = os.path.join(PYMC_FIGS, 'trace_plots.png')
    if os.path.exists(chart):
        slide.shapes.add_picture(chart, Inches(0.40), Inches(1.58),
                                 Inches(12.40), Inches(5.20))
    _add_textbox(slide,
                 'R-hat max = 1.0 (pass < 1.1)  ·  ESS min = 578 (pass > 400)  ·  '
                 'Chains well-mixed — no divergences',
                 left=0.50, top=6.85, width=12.20, height=0.40,
                 font_size=11, color=CHARCOAL)


# ── A3: Channel Taxonomy ───────────────────────────────────────────────────────
TAXONOMY = [
    ('facebook',    'Meta',                  'iOS / Android / Web'),
    ('google',      'Google Search',         'iOS / Android / Web'),
    ('tiktok',      'TikTok',                'iOS / Android / Web'),
    ('applovin',    'DSP: Applovin',         'Separate'),
    ('liftoff',     'DSP: Liftoff',          'Separate'),
    ('Inmobidsp',   'DSP: InMobi',           'Separate'),
    ('Stackadapt +\nAppvertiser', 'DSP: Combined', 'Grouped (D004)'),
    ('linear_tv',   'Linear TV',             'Separate'),
    ('ctv',         'CTV',                   'Separate'),
    ('podcast',     'Podcast / Audio',       'Separate'),
    ('apple_search','Apple Search Ads',      'Separate'),
    ('influencer',  'Influencers',           'Separate'),
    ('bing',        '(excluded)',             'Dropped — too small'),
    ('Others',      '(excluded)',             'Erroneous data removed'),
]


def _a3(prs):
    slide = add_content_slide(prs, 'Appendix A3 — Channel Taxonomy (13 Modeled Channels)',
                              page_num='A3')
    _add_rect(slide, 0.50, 1.62, 12.20, 0.38, fill_rgb=NAVY)
    for h, lx in [('SOURCE_GROUP', 0.60), ('Canonical Channel', 3.00), ('Model Treatment', 7.80)]:
        _add_textbox(slide, h, left=lx, top=1.66, width=4.00, height=0.28,
                     font_size=11, bold=True, color=WHITE)

    for ri, (sg, ch, tr) in enumerate(TAXONOMY):
        rt = 2.02 + ri * 0.35
        from pptx.dml.color import RGBColor as _RGB
        if ri % 2 == 1:
            _add_rect(slide, 0.50, rt, 12.20, 0.35, fill_rgb=_RGB(0xF0, 0xF0, 0xF0))
        from pptx.dml.color import RGBColor as _RGB
        for val, lx in [(sg, 0.60), (ch, 3.00), (tr, 7.80)]:
            tc = CHARCOAL if 'excluded' not in ch else _RGB(0xAA, 0xAA, 0xAA)
            _add_textbox(slide, val, left=lx, top=rt + 0.04, width=4.00, height=0.28,
                         font_size=10, color=tc)


# ── A4: Incrementality Tests ───────────────────────────────────────────────────
INCR = [
    ('TikTok',      'Aug–Sep 2025', 'User-level holdout (high conviction)',
     'iOS $108.83 / Android $81.68 / Web $112.12'),
    ('Meta May 2025','May 2025',    '3-cell design (reach restricted)',
     'iOS $135.48 / Android $63.06 / Web $156.89'),
    ('Meta Jan 2026','Jan 2026',    'Cancelled — use with wide CI',
     'iOS $1,350 / Android $154.52 / Web $239.75'),
    ('CTV',         'Oct–Nov 2025', 'Geo holdout (lower conviction)',
     'All platforms: $135'),
]


def _a4(prs):
    slide = add_content_slide(prs, 'Appendix A4 — Incrementality Tests Used as Model Priors',
                              page_num='A4')
    _add_rect(slide, 0.50, 1.62, 12.20, 0.38, fill_rgb=NAVY)
    for h, lx in [('Channel', 0.60), ('Window', 2.70), ('Test Design', 4.50), ('iCAC Values', 8.60)]:
        _add_textbox(slide, h, left=lx, top=1.66, width=3.50, height=0.28,
                     font_size=11, bold=True, color=WHITE)

    for ri, (ch, win, design, icac) in enumerate(INCR):
        rt = 2.10 + ri * 0.90
        from pptx.dml.color import RGBColor as _RGB
        if ri % 2 == 1:
            _add_rect(slide, 0.50, rt, 12.20, 0.88, fill_rgb=_RGB(0xF0, 0xF0, 0xF0))
        for val, lx, w in [(ch, 0.60, 2.00), (win, 2.70, 1.70),
                            (design, 4.50, 4.00), (icac, 8.60, 4.00)]:
            _add_textbox(slide, val, left=lx, top=rt + 0.06, width=w, height=0.75,
                         font_size=11, color=CHARCOAL)


# ── A5: Channel Contributions ─────────────────────────────────────────────────
def _a5(prs):
    slide = add_content_slide(prs, 'Appendix A5 — Channel Contributions (Facebook POC)',
                              page_num='A5')
    chart = os.path.join(PYMC_FIGS, 'channel_contributions.png')
    if os.path.exists(chart):
        slide.shapes.add_picture(chart, Inches(0.40), Inches(1.58),
                                 Inches(12.40), Inches(5.20))


# ── A6: EDA Scatter ───────────────────────────────────────────────────────────
def _a6(prs):
    slide = add_content_slide(prs, 'Appendix A6 — Spend vs. Conversions Scatter',
                              page_num='A6')
    chart = os.path.join(EDA_FIGS, 'spend_vs_conversions_scatter.png')
    if os.path.exists(chart):
        slide.shapes.add_picture(chart, Inches(0.40), Inches(1.58),
                                 Inches(12.40), Inches(5.20))


def build(prs):
    # Section divider
    from templates import add_title_section_slide
    add_title_section_slide(prs, title='Appendix',
                            subtitle='Supporting charts, tables, and model diagnostics')
    _a1(prs)
    _a2(prs)
    _a3(prs)
    _a4(prs)
    _a5(prs)
    _a6(prs)
