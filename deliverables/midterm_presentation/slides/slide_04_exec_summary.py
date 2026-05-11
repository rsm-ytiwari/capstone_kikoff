import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from templates import add_content_slide, _add_rect, _add_textbox, NAVY, GOLD, WHITE, CHARCOAL
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

ASSET = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                     'assets', 'slide_04_exec_summary_icons.png')

COLS = [
    {
        'header': 'The Problem',
        'color': NAVY,
        'bullets': [
            '$113M annual spend across 13+ channels',
            'Platform-reported ROAS double-counts and ignores baseline',
            'No causal attribution → budget allocation is guesswork',
        ],
    },
    {
        'header': 'Our Approach',
        'color': NAVY,
        'bullets': [
            'Bayesian Media Mix Model (PyMC-Marketing)',
            '4 real incrementality tests as calibration priors',
            '639 daily observations, 2024–2026',
        ],
    },
    {
        'header': 'Early Findings',
        'color': GOLD,
        'bullets': [
            'Facebook POC: model converged (R-hat = 1.0)',
            'Facebook-iOS: $34 iCAC, improving — efficient',
            'Facebook-Web: $234 iCAC, declining — hold spend',
        ],
    },
]


def _column(slide, col, left, top=1.70, width=3.80, height=5.00):
    # Header
    _add_rect(slide, left, top, width, 0.07, fill_rgb=col['color'])
    _add_textbox(slide, col['header'],
                 left=left, top=top + 0.10, width=width, height=0.55,
                 font_size=17, bold=True, color=NAVY)
    # Bullets
    body_top = top + 0.75
    for bullet in col['bullets']:
        txb = slide.shapes.add_textbox(
            Inches(left), Inches(body_top), Inches(width), Inches(0.90)
        )
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = '• ' + bullet
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.color.rgb = CHARCOAL
        run.font.name = 'Calibri'
        body_top += 1.20


def build(prs):
    slide = add_content_slide(
        prs,
        'Phase 1 Complete, POC Converged — Facebook-iOS Is Efficient; Web and Android Are Not',
        page_num=4
    )

    if os.path.exists(ASSET):
        slide.shapes.add_picture(ASSET, Inches(0.40), Inches(1.55),
                                 Inches(12.40), Inches(5.30))
    else:
        col_starts = [0.50, 4.60, 8.70]
        for col, left in zip(COLS, col_starts):
            _column(slide, col, left)
