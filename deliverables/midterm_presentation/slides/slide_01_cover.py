import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from templates import add_title_section_slide, _add_textbox, GOLD, WHITE, CHARCOAL
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def build(prs):
    slide = add_title_section_slide(
        prs,
        title='Marketing Mix Modeling\nfor Kikoff',
        subtitle='Rady School of Management  ·  MSBA Capstone  ·  2026',
    )

    # Team names
    team = [
        'Chunjiang Liu  |  Shivang Bhatt  |  Yash Tiwari  |  Yoko He',
    ]
    txb = slide.shapes.add_textbox(
        Inches(1.00), Inches(5.80), Inches(9.50), Inches(0.70)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = team[0]
    run.font.size = Pt(15)
    run.font.bold = False
    run.font.color.rgb = WHITE
    run.font.name = 'Calibri'

    # Sponsor line
    txb2 = slide.shapes.add_textbox(
        Inches(1.00), Inches(6.50), Inches(9.50), Inches(0.60)
    )
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = 'Sponsor: Abheek Sinha, Kikoff  ·  Supervisor: Keith'
    run2.font.size = Pt(13)
    run2.font.bold = False
    run2.font.color.rgb = GOLD
    run2.font.name = 'Calibri'
