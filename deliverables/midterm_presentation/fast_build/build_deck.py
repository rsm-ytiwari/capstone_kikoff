"""
Fast Build — Kikoff MMM Midterm Deck (Run B)
All writes isolated to: deliverables/midterm_presentation/fast_build/

Run from project root:
    my-notebook-project/.venv/bin/python deliverables/midterm_presentation/fast_build/build_deck.py
"""
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────────────────
_HERE       = os.path.dirname(os.path.abspath(__file__))
ROOT        = os.path.abspath(os.path.join(_HERE, "..", "..", ".."))
ASSETS_SRC  = os.path.join(ROOT, "deliverables", "midterm_presentation", "assets")
FAST_ASSETS = os.path.join(_HERE, "assets")
EDA_FIGS    = os.path.join(ROOT, "outputs", "P1_02_eda", "figures")
PYMC_FIGS   = os.path.join(ROOT, "outputs", "P2_01_pymc", "figures")
OUTPUT      = os.path.join(_HERE, "kikoff_mmm_progress_review.pptx")

BG_CONTENT   = os.path.join(ASSETS_SRC, "content.png")
BG_TITLE     = os.path.join(ASSETS_SRC, "main_background.png")
TEAM_IMG     = os.path.join(ASSETS_SRC, "Draft_MMM_kikoff.pptx.png")

os.makedirs(FAST_ASSETS, exist_ok=True)

# ── Brand colors ───────────────────────────────────────────────────────────────
NAVY     = RGBColor(0x18, 0x2B, 0x48)
GOLD     = RGBColor(0xFF, 0xCC, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x33, 0x33, 0x33)
LGRAY    = RGBColor(0xF0, 0xF0, 0xF0)
RED_RGB  = RGBColor(0xCC, 0x00, 0x00)

NAVY_HEX     = "#182B48"
GOLD_HEX     = "#FFCC00"
WHITE_HEX    = "#FFFFFF"
CHARCOAL_HEX = "#333333"
LGRAY_HEX    = "#F0F0F0"


# ══════════════════════════════════════════════════════════════════════════════
# Slide / text helpers
# ══════════════════════════════════════════════════════════════════════════════

def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _txt(slide, left, top, w, h, text, size, bold=False, color=None,
         align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tb.fill.background()  # prevent white-box artifact in non-PowerPoint renderers
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name   = "Calibri"
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color if color else CHARCOAL
    return tb


def _bullets(slide, left, top, w, h, items, size=15, color=None):
    """items: list of str or (str, bool) for (text, bold)."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        text, bold = item if isinstance(item, tuple) else (item, False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = text
        r.font.name  = "Calibri"
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = color if color else CHARCOAL
        first = False
    return tb


def _note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _rect(slide, left, top, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        1,  # rectangle
        Inches(left), Inches(top), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    bc = border_color if border_color else fill_color
    shape.line.color.rgb = bc
    return shape


def _img(slide, path, left, top, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))
    else:
        _txt(slide, left, top, w, h, f"[image not found: {os.path.basename(path)}]",
             11, color=RED_RGB)


# ── Frame builders ─────────────────────────────────────────────────────────────

def _title_slide(prs, title, subtitle=None, detail=None):
    slide = _blank_slide(prs)
    _img(slide, BG_TITLE, 0, 0, 13.333, 7.5)
    _txt(slide, 1.00, 2.00, 9.50, 2.00, title, 34, bold=True, color=WHITE)
    if subtitle:
        _txt(slide, 1.00, 4.10, 9.50, 0.70, subtitle, 20, color=GOLD)
    if detail:
        _txt(slide, 1.00, 4.90, 9.50, 1.80, detail, 16, color=WHITE)
    return slide


def _content_slide(prs, page_num, title):
    slide = _blank_slide(prs)
    _img(slide, BG_CONTENT, 0, 0, 13.333, 7.5)
    # Cover baked-in "2" with navy rectangle then actual page number
    _rect(slide, 0.13, 7.12, 0.34, 0.34, NAVY)
    _txt(slide, 0.14, 7.13, 0.32, 0.30, str(page_num), 9, color=GOLD)
    # Slide title
    _txt(slide, 0.50, 0.87, 9.50, 0.65, title, 21, bold=True, color=NAVY)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# Asset generators (matplotlib → fast_build/assets/)
# ══════════════════════════════════════════════════════════════════════════════

def _gen_problem_diagram():
    """Redesigned: clean L-shaped fan-in connectors, larger text, clear panels."""
    path = os.path.join(FAST_ASSETS, "slide_05_problem_diagram.png")

    fig = plt.figure(figsize=(13, 4.8), facecolor="white")

    # Two side-by-side axes
    ax_l = fig.add_axes([0.02, 0.08, 0.44, 0.86])
    ax_r = fig.add_axes([0.54, 0.08, 0.44, 0.86])

    for ax in (ax_l, ax_r):
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 4.8)
        ax.set_facecolor("white")
        ax.axis("off")

    # Vertical divider between panels
    fig.add_artist(plt.Line2D(
        [0.495, 0.495], [0.06, 0.96],
        transform=fig.transFigure,
        color=LGRAY_HEX, lw=2, linestyle="--"
    ))

    channel_ys = [3.6, 2.2, 0.8]   # y-centers for Meta, TikTok, Google
    junction_x = 6.0                # fan-in x
    outcome_x  = 7.0                # outcome box left edge
    mid_y      = 2.2                # arrow midpoint y

    # ── Left panel: Current State ──────────────────────────────────────────────
    ax_l.text(5.0, 4.55, "Current State", fontsize=15, fontweight="bold",
              color=NAVY_HEX, ha="center")

    channels_l = ["Meta", "TikTok", "Google"]
    for label, cy in zip(channels_l, channel_ys):
        r = FancyBboxPatch((0.2, cy - 0.48), 4.6, 0.96,
                           boxstyle="round,pad=0.10",
                           facecolor=LGRAY_HEX, edgecolor=NAVY_HEX, lw=1.8)
        ax_l.add_patch(r)
        ax_l.text(2.5, cy, f"{label}  (self-reported)",
                  fontsize=12, ha="center", va="center", color=CHARCOAL_HEX)
        # L-shaped connector: horizontal from box right → junction_x
        ax_l.plot([4.82, junction_x], [cy, cy], color=CHARCOAL_HEX, lw=1.6)
        # Vertical from cy to mid_y (only if not already at mid_y)
        if abs(cy - mid_y) > 0.05:
            ax_l.plot([junction_x, junction_x], [cy, mid_y],
                      color=CHARCOAL_HEX, lw=1.6)

    # Arrow from junction to outcome
    ax_l.annotate("", xy=(outcome_x, mid_y), xytext=(junction_x, mid_y),
                  arrowprops=dict(arrowstyle="->", color="#CC0000", lw=2.0,
                                  mutation_scale=18))

    # Bad outcome box
    bad = FancyBboxPatch((outcome_x, mid_y - 0.80), 2.8, 1.6,
                         boxstyle="round,pad=0.10",
                         facecolor="#FFE0E0", edgecolor="#CC0000", lw=2.0)
    ax_l.add_patch(bad)
    ax_l.text(outcome_x + 1.4, mid_y,
              "Budget decisions\nwithout causal\nattribution",
              fontsize=10, ha="center", va="center", color="#CC0000",
              multialignment="center")

    ax_l.text(5.0, 0.08,
              "Each platform overclaims its credit",
              fontsize=10, ha="center", color="#CC0000", style="italic")

    # ── Right panel: MMM State ─────────────────────────────────────────────────
    ax_r.text(5.0, 4.55, "With Bayesian MMM", fontsize=15, fontweight="bold",
              color=NAVY_HEX, ha="center")

    channels_r = ["Meta (iOS / Web)", "TikTok (iOS / Web)", "Google (iOS / Web)"]
    for label, cy in zip(channels_r, channel_ys):
        r = FancyBboxPatch((0.2, cy - 0.48), 4.6, 0.96,
                           boxstyle="round,pad=0.10",
                           facecolor=LGRAY_HEX, edgecolor=NAVY_HEX, lw=1.8)
        ax_r.add_patch(r)
        ax_r.text(2.5, cy, label,
                  fontsize=12, ha="center", va="center", color=CHARCOAL_HEX)
        ax_r.plot([4.82, junction_x], [cy, cy], color=CHARCOAL_HEX, lw=1.6)
        if abs(cy - mid_y) > 0.05:
            ax_r.plot([junction_x, junction_x], [cy, mid_y],
                      color=CHARCOAL_HEX, lw=1.6)

    ax_r.annotate("", xy=(outcome_x, mid_y), xytext=(junction_x, mid_y),
                  arrowprops=dict(arrowstyle="->", color="#2E7D32", lw=2.0,
                                  mutation_scale=18))

    good = FancyBboxPatch((outcome_x, mid_y - 0.80), 2.8, 1.6,
                          boxstyle="round,pad=0.10",
                          facecolor="#E8F5E9", edgecolor="#2E7D32", lw=2.0)
    ax_r.add_patch(good)
    ax_r.text(outcome_x + 1.4, mid_y,
              "Causal ROAS\nper channel\nx platform",
              fontsize=10, ha="center", va="center", color="#2E7D32",
              multialignment="center")

    ax_r.text(5.0, 0.08,
              "Incrementality priors time-limited to test windows",
              fontsize=9.5, ha="center", color="#2E7D32", style="italic")

    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  Generated: {path}")
    return path


def _gen_framework_table():
    path = os.path.join(FAST_ASSETS, "slide_08_framework_table.png")

    rows = [
        ["PyMC-Marketing", "Full posterior", "Native (add_lift_test)", "Yes", ">> SELECTED"],
        ["Meridian (Google)", "Full posterior", "Limited",              "No",  "Rejected: less flexible"],
        ["LightweightMMM",   "Full posterior", "Manual workaround",    "No",  "Rejected: deprecated"],
        ["Robyn (Meta)",     "Bootstrap CI",   "Manual workaround",    "No",  "Rejected: no uncertainty"],
        ["Custom PyMC",      "Full posterior", "Manual workaround",    "Yes", "Rejected: slower dev"],
    ]
    cols = ["Framework", "Uncertainty", "Lift Integration", "Dual DV", "Decision"]

    fig, ax = plt.subplots(figsize=(12, 4.0), facecolor="white")
    ax.axis("off")

    tbl = ax.table(cellText=rows, colLabels=cols, cellLoc="center", loc="center")
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(11.5)
    tbl.scale(1, 2.3)

    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor(LGRAY_HEX)
        cell.set_linewidth(0.5)
        if r == 0:
            cell.set_facecolor(NAVY_HEX)
            cell.set_text_props(color="white", fontweight="bold")
        elif r == 1:  # winner row
            cell.set_facecolor("#E8F5E9")
            cell.set_text_props(color=CHARCOAL_HEX)
        else:
            cell.set_facecolor("white")
            cell.set_text_props(color=CHARCOAL_HEX)

    fig.tight_layout()
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  Generated: {path}")
    return path


def _gen_appendix_phases():
    path = os.path.join(FAST_ASSETS, "appendix_phases.png")

    rows = [
        ["Phase 1\nBusiness & Data Audit",
         "Lock channel taxonomy (15 source groups → 13)\nSelect modeling framework\nClean 21 months of daily data",
         "13 channels locked · PyMC-Marketing approved\n5 data quality issues resolved · 639-row clean dataset",
         "COMPLETE"],
        ["Phase 2\nModel Build & POC",
         "Fit Bayesian MMM on Meta channels (iOS, Android, Web)\nValidate convergence · produce first ICAC / iROAS curves",
         "Meta × 3 platforms: ICAC + iROAS curves\nwith posterior uncertainty bands",
         "IN PROGRESS"],
        ["Phase 3\nChannel Replication",
         "Run validated model for all 13 channel groups × platform\nBuild decisioning layer (scale / hold / investigate / test)",
         "Full ICAC / iROAS table · saturation reads\nSeasonal budget scenarios",
         "NOT STARTED"],
        ["Phase 4\nDeliverables & Handoff",
         "Seasonal scenarios (tax season, holiday, steady-state)\nFinal presentation · Reproducible model package",
         "June 6 final presentation\nHandoff package to Kikoff",
         "NOT STARTED"],
    ]
    cols = ["Phase", "Work Scope", "Key Outputs", "Status"]

    fig, ax = plt.subplots(figsize=(13, 6.0), facecolor="white")
    ax.axis("off")

    tbl = ax.table(cellText=rows, colLabels=cols, cellLoc="left", loc="center")
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(9.5)
    tbl.scale(1, 4.8)

    GREEN_HEX  = "#E8F5E9"
    YELLOW_HEX = "#FFF9C4"
    STATUS_COLORS = {"COMPLETE": "#2E7D32", "IN PROGRESS": NAVY_HEX,
                     "NOT STARTED": "#888888"}
    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor(LGRAY_HEX)
        cell.set_linewidth(0.8)
        cell.PAD = 0.06
        if r == 0:
            cell.set_facecolor(NAVY_HEX)
            cell.set_text_props(color="white", fontweight="bold", ha="center")
        elif r == 1:
            cell.set_facecolor(GREEN_HEX)
            txt_color = STATUS_COLORS.get(rows[r-1][3], CHARCOAL_HEX) if c == 3 else CHARCOAL_HEX
            cell.set_text_props(color=txt_color, fontweight="bold" if c == 3 else "normal")
        elif r == 2:
            cell.set_facecolor(YELLOW_HEX)
            txt_color = STATUS_COLORS.get(rows[r-1][3], CHARCOAL_HEX) if c == 3 else CHARCOAL_HEX
            cell.set_text_props(color=txt_color, fontweight="bold" if c == 3 else "normal")
        else:
            cell.set_facecolor("white")
            txt_color = STATUS_COLORS.get(rows[r-1][3], CHARCOAL_HEX) if c == 3 else CHARCOAL_HEX
            cell.set_text_props(color=txt_color, fontweight="bold" if c == 3 else "normal")

    col_widths = [0.14, 0.38, 0.36, 0.12]
    for col_idx, w in enumerate(col_widths):
        for row_idx in range(len(rows) + 1):
            tbl[row_idx, col_idx].set_width(w)

    fig.tight_layout(pad=0.3)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  Generated: {path}")
    return path


def _gen_appendix_milestones():
    path = os.path.join(FAST_ASSETS, "appendix_milestones.png")

    rows = [
        ["M1  Data Cleaning",
         "LTV imputation · platform label normalization · DSP channel mapping",
         "5 gate checks pass (z-score, row count, no nulls)",
         "COMPLETE"],
        ["M2  Time-Limited Priors",
         "Choose how to apply incrementality test results only within test windows\n(not retroactively across full history)",
         "Mechanism selected (add_lift_test_measurements)\nD007 = testing strategy decision (not final Q19 resolution)",
         "COMPLETE"],
        ["M3  Meta POC Model",
         "Fit PyMC-Marketing on Meta iOS, Android, Web\nConfirm model health and produce first curves",
         "R-hat < 1.1  ·  ESS > 400  ·  ICAC in $30–$500 range\niCAC + iROAS curves in outputs/",
         "IN PROGRESS"],
        ["M4  Others Cleanup",
         "Confirm with sponsor that Idle/Others erroneous data removal\nwas intentional · close open question Q18",
         "Sponsor confirmation received",
         "NON-BLOCKING"],
        ["M5  Decisioning Layer",
         "Build per-channel table: spend signal, iCAC/iROAS trend,\nsaturation read, recommended action, spend-move %",
         "Template populated for Meta channels",
         "AFTER M3"],
        ["M6  Full Replication",
         "Run M3 + M5 for all remaining channel groups × platform\n(TikTok, Google, CTV, DSPs, and tail channels)",
         "Complete deliverable package — all 13 channels",
         "FINAL STEP"],
    ]
    cols = ["Milestone", "What It Means", "Done When", "Status"]

    fig, ax = plt.subplots(figsize=(13, 6.5), facecolor="white")
    ax.axis("off")

    tbl = ax.table(cellText=rows, colLabels=cols, cellLoc="left", loc="center")
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(9.5)
    tbl.scale(1, 4.5)

    GREEN_HEX  = "#E8F5E9"
    YELLOW_HEX = "#FFF9C4"
    STATUS_COLORS = {"COMPLETE": "#2E7D32", "IN PROGRESS": NAVY_HEX,
                     "NON-BLOCKING": "#888888", "AFTER M3": "#888888",
                     "FINAL STEP": "#888888"}
    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor(LGRAY_HEX)
        cell.set_linewidth(0.8)
        cell.PAD = 0.06
        if r == 0:
            cell.set_facecolor(NAVY_HEX)
            cell.set_text_props(color="white", fontweight="bold", ha="center")
        elif r in (1, 2):
            cell.set_facecolor(GREEN_HEX)
            txt_color = STATUS_COLORS.get(rows[r-1][3], CHARCOAL_HEX) if c == 3 else CHARCOAL_HEX
            cell.set_text_props(color=txt_color, fontweight="bold" if c == 3 else "normal")
        elif r == 3:
            cell.set_facecolor(YELLOW_HEX)
            txt_color = STATUS_COLORS.get(rows[r-1][3], CHARCOAL_HEX) if c == 3 else CHARCOAL_HEX
            cell.set_text_props(color=txt_color, fontweight="bold" if c == 3 else "normal")
        else:
            cell.set_facecolor("white")
            txt_color = STATUS_COLORS.get(rows[r-1][3], CHARCOAL_HEX) if c == 3 else CHARCOAL_HEX
            cell.set_text_props(color=txt_color, fontweight="bold" if c == 3 else "normal")

    col_widths = [0.14, 0.42, 0.32, 0.12]
    for col_idx, w in enumerate(col_widths):
        for row_idx in range(len(rows) + 1):
            tbl[row_idx, col_idx].set_width(w)

    fig.tight_layout(pad=0.3)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  Generated: {path}")
    return path


def _gen_phases():
    path = os.path.join(FAST_ASSETS, "slide_09_phases.png")
    fig, ax = plt.subplots(figsize=(12, 3.2), facecolor="white")
    ax.set_xlim(-0.3, 6.0)
    ax.set_ylim(0, 3.2)
    ax.axis("off")

    phases = [
        ("Phase 1\nBusiness &\nData Audit",   "COMPLETE",    GOLD_HEX,    NAVY_HEX,  True),
        ("Phase 2\nModel Build\n& POC",        "IN PROGRESS", "#FFF9C4",   NAVY_HEX,  True),
        ("Phase 3\nChannel\nReplication",      "~Weeks 3-4",  "white",     LGRAY_HEX, False),
        ("Phase 4\nScenario\nPlanning",        "~Week 5",     "white",     LGRAY_HEX, False),
        ("Phase 5\nSynthesis\n& Handoff",      "~Week 6",     "white",     LGRAY_HEX, False),
    ]

    bw, bh, gap = 1.05, 1.9, 0.10
    step = bw + gap

    for i, (label, status, fc, ec, active) in enumerate(phases):
        x = i * step
        rect = FancyBboxPatch((x, 0.6), bw, bh, boxstyle="round,pad=0.06",
                              facecolor=fc, edgecolor=ec, linewidth=2.0)
        ax.add_patch(rect)
        ax.text(x + bw / 2, 1.55, label, fontsize=9.5, ha="center", va="center",
                color=NAVY_HEX, fontweight="bold" if active else "normal",
                multialignment="center", linespacing=1.4)
        ax.text(x + bw / 2, 0.75, status, fontsize=8.5, ha="center", va="center",
                color=NAVY_HEX if active else CHARCOAL_HEX, style="italic")

        if i < 4:
            ax.annotate("", xy=(x + bw + gap, 1.55), xytext=(x + bw, 1.55),
                        arrowprops=dict(arrowstyle="->", color=CHARCOAL_HEX, lw=1.5))

    # "We are here" annotation — Phase 2 (i=1)
    px = step + bw / 2
    ax.annotate(
        "We are here",
        xy=(px, 2.52), xytext=(px + 0.8, 3.0),
        fontsize=10, color=NAVY_HEX, fontweight="bold",
        arrowprops=dict(arrowstyle="->", color=NAVY_HEX, lw=1.5),
    )

    fig.tight_layout(pad=0.4)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  Generated: {path}")
    return path


def _gen_deliverable_preview():
    import numpy as np
    path = os.path.join(FAST_ASSETS, "slide_07_deliverable_preview.png")

    fig, (ax_curve, ax_phase) = plt.subplots(
        1, 2, figsize=(8.5, 4.6), facecolor="white",
        gridspec_kw={"width_ratios": [1.35, 1], "wspace": 0.06},
    )
    fig.subplots_adjust(left=0.07, right=0.97, top=0.88, bottom=0.18)

    # ── Panel 1: ICAC Diminishing-Returns Curve ───────────────────────────────
    spend = np.linspace(2_000, 22_000, 300)

    def icac_fn(s):
        return 60 * (s / 1000) ** 0.35  # illustrative power law

    icac      = icac_fn(spend)
    spend_mid = 12_000
    spend_lo  = spend_mid * 0.8
    spend_hi  = spend_mid * 1.2

    ax_curve.axvspan(spend_lo, spend_hi, alpha=0.13, color=GOLD_HEX, zorder=0)
    ax_curve.axvline(spend_mid, color=CHARCOAL_HEX, lw=0.8, ls="--",
                     alpha=0.35, zorder=1)
    ax_curve.plot(spend, icac, color=NAVY_HEX, lw=2.5, zorder=2)

    xs = [spend_lo, spend_mid, spend_hi]
    ys = [icac_fn(x) for x in xs]
    ax_curve.scatter(xs, ys, s=[55, 90, 55], color=GOLD_HEX,
                     edgecolors=NAVY_HEX, linewidths=1.2, zorder=4)

    for sx, sy, lbl, dx in zip(
        xs, ys,
        ["−20%", "Current", "+20%"],
        [-700, 0, 700],
    ):
        ax_curve.text(sx + dx, sy + 8, lbl,
                      fontsize=8, color=NAVY_HEX, fontweight="bold",
                      ha="center", va="bottom")

    ax_curve.text(0.5, 0.44, "Illustrative",
                  transform=ax_curve.transAxes,
                  fontsize=26, color="gray", alpha=0.07,
                  ha="center", va="center", rotation=28, fontweight="bold")

    ax_curve.spines["top"].set_visible(False)
    ax_curve.spines["right"].set_visible(False)
    ax_curve.spines["left"].set_color(LGRAY_HEX)
    ax_curve.spines["bottom"].set_color(LGRAY_HEX)
    ax_curve.tick_params(labelbottom=False, labelleft=False, length=0)
    ax_curve.set_xlabel("Weekly Spend  →", fontsize=8.5, color=CHARCOAL_HEX)
    ax_curve.set_ylabel("Incremental CAC ($)  →", fontsize=8.5, color=CHARCOAL_HEX)
    ax_curve.set_title("ICAC Curve — per channel × platform",
                       fontsize=9.5, color=NAVY_HEX, fontweight="bold", pad=5)

    # ── Panel 2: Channel Phase Matrix ─────────────────────────────────────────
    ax_phase.axis("off")
    ax_phase.set_xlim(0, 1)
    ax_phase.set_ylim(0, 1)

    ax_phase.add_patch(FancyBboxPatch(
        (0.02, 0.83), 0.96, 0.14, boxstyle="round,pad=0.01",
        facecolor=NAVY_HEX, edgecolor="none"))
    ax_phase.text(0.50, 0.905, "Channel Coverage",
                  ha="center", va="center", fontsize=9.5,
                  color="white", fontweight="bold")

    ax_phase.add_patch(FancyBboxPatch(
        (0.02, 0.47), 0.96, 0.32, boxstyle="round,pad=0.01",
        facecolor="#FFF9C4", edgecolor=LGRAY_HEX, linewidth=0.7))
    ax_phase.text(0.08, 0.74, "Phase 2",
                  fontsize=9, fontweight="bold", color=NAVY_HEX, va="center")
    ax_phase.text(0.08, 0.59, "meta_web · meta_ios\nmeta_android",
                  fontsize=8, color=CHARCOAL_HEX, va="center")

    ax_phase.add_patch(FancyBboxPatch(
        (0.02, 0.11), 0.96, 0.32, boxstyle="round,pad=0.01",
        facecolor="#E3F2FD", edgecolor=LGRAY_HEX, linewidth=0.7))
    ax_phase.text(0.08, 0.38, "Phase 3",
                  fontsize=9, fontweight="bold", color="#1565C0", va="center")
    ax_phase.text(0.08, 0.23, "tiktok_ios · tiktok_android\ngoogle / ctv / …",
                  fontsize=8, color=CHARCOAL_HEX, va="center")

    ax_phase.set_title("Phase Rollout",
                       fontsize=9.5, color=NAVY_HEX, fontweight="bold", pad=5)

    # ── Footer ─────────────────────────────────────────────────────────────────
    fig.text(0.5, 0.04,
             "13 channels × platform — results at final presentation, June 6",
             ha="center", fontsize=8.5, color=NAVY_HEX, style="italic")

    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  Generated: {path}")
    return path


def _gen_milestones():
    path = os.path.join(FAST_ASSETS, "slide_12_milestones.png")
    fig, ax = plt.subplots(figsize=(12, 4.2), facecolor="white")
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 4.2)
    ax.axis("off")

    rows = [
        (
            "[DONE]  Phase 1 Complete",
            "Taxonomy locked  *  PyMC-Marketing approved  *  5/5 M1 data checks passed  *  Clean dataset in outputs/",
            "#E8F5E9", "#2E7D32",
        ),
        (
            "[NOW]  Next 2-3 Weeks",
            "Mechanism 2 (time-windowed priors) implementation  *  M3: Meta POC health check (R-hat < 1.1, ESS > 400, ICAC in range)  *  M5: Decisioning layer",
            "#FFF9C4", NAVY_HEX,
        ),
        (
            "[IMPACT]  Downstream",
            "Defensible, repeatable model anchoring Kikoff's next $100M in marketing spend -- clear saturation reads + test prioritization",
            "#E3F2FD", "#1565C0",
        ),
    ]

    row_h = 1.18
    for idx, (header, body, fc, tc) in enumerate(rows):
        y_top = 3.95 - idx * row_h * 1.05
        r, g, b = int(fc[1:3], 16), int(fc[3:5], 16), int(fc[5:7], 16)
        rect = FancyBboxPatch((0.2, y_top - row_h + 0.08), 11.6, row_h - 0.08,
                              boxstyle="round,pad=0.08",
                              facecolor=(r / 255, g / 255, b / 255),
                              edgecolor=LGRAY_HEX, linewidth=1.2)
        ax.add_patch(rect)
        ax.text(0.55, y_top - 0.30, header, fontsize=13, fontweight="bold",
                color=tc, va="center")
        ax.text(0.55, y_top - 0.75, body, fontsize=10.5, color=CHARCOAL_HEX,
                va="center", wrap=True)

    fig.tight_layout(pad=0.3)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  Generated: {path}")
    return path


# ══════════════════════════════════════════════════════════════════════════════
# Slide builders
# ══════════════════════════════════════════════════════════════════════════════

def slide_01(prs):
    slide = _title_slide(
        prs,
        title="Measuring What Actually Drives\nKikoff's Growth",
        subtitle="UC San Diego Rady MSBA — Capstone 2026",
        detail="Yash Tiwari  ·  Chunjiang Liu  ·  Shivang Bhatt  ·  Yoko He\nMay 2026",
    )
    _note(slide, (
        "Hi everyone — thank you for being here. I'm Yash Tiwari, and I'm here with my teammates "
        "Chunjiang Liu, Shivang Bhatt, and Yoko He. We're the Kikoff capstone team from the UC San Diego "
        "Rady MSBA program, and today we'll be walking you through our progress on building a Media Mix Model for Kikoff."
    ))


def slide_02(prs):
    slide = _content_slide(prs, 2, "Four Sections, One Story: From Data to Budget Decisions")
    sections = [
        ("1.  Business Context & Data Foundation",
         "Who is Kikoff, what problem we're solving, and what data we're working with"),
        ("2.  Methodology & Deliverables",
         "What we're building, how we selected our framework, and what we evaluated"),
        ("3.  Findings to Date",
         "Key EDA results, data quality resolutions, and early modeling signals"),
        ("4.  Next Steps",
         "Three concrete milestones between now and the final deliverable"),
    ]
    y = 1.70
    for header, sub in sections:
        _txt(slide, 0.75, y, 11.80, 0.42, header, 17, bold=True, color=NAVY)
        _txt(slide, 1.10, y + 0.42, 11.40, 0.40, sub, 14, color=CHARCOAL)
        y += 1.08
    _note(slide, (
        "Here's a quick roadmap of what we'll cover. We'll move from the business context and data "
        "foundation, through our methodology and deliverables, into our findings to date, and close "
        "with recommendations and next steps. Each of us will take a section."
    ))


def slide_03(prs):
    """Team slide — embed Draft_MMM_kikoff.pptx.png as the main content visual."""
    slide = _content_slide(prs, 3, "Meet the Team")
    _img(slide, TEAM_IMG, 0.38, 1.52, 12.55, 5.62)
    _note(slide, (
        "A quick introduction. [Coordinate live — each person introduces themselves in one sentence "
        "as they're named: name, concentration or prior role, and what they owned on this project.]"
    ))


def slide_04(prs):
    slide = _content_slide(prs, 4,
                           "Platform Attribution Costs Kikoff Precision on $112M in Annual Spend")

    CARD_TOP = 1.62
    CARD_H   = 5.30
    BAND_H   = 0.56   # colored header band
    cw       = 3.92
    gap      = 0.24

    cards = [
        (
            "The Problem",
            [
                "$112M annual spend — zero causal attribution",
                "Each platform self-reports ROAS — all overclaim",
                "Last-touch credits the final click, not the cause",
            ],
            "#E3F2FD", "#1565C0", WHITE,
        ),
        (
            "Key Pivots",
            [
                "Channel  →  channel × platform",
                "Single DV  →  conversions + LTV",
                "Linear regression  →  Bayesian with priors",
            ],
            "#FFF9C4", "#E65100", WHITE,
        ),
        (
            "Where We Stand",
            [
                "✓  Framework: PyMC-Marketing — approved",
                "✓  13 channels locked, taxonomy clean",
                "✓  Phase 1 complete — all 5 gates passed",
            ],
            "#E8F5E9", "#2E7D32", WHITE,
        ),
    ]

    for i, (header, bullets, card_fc, band_fc, band_tc) in enumerate(cards):
        lx = 0.38 + i * (cw + gap)
        # Card background
        r, g, b = int(card_fc[1:3], 16), int(card_fc[3:5], 16), int(card_fc[5:7], 16)
        _rect(slide, lx, CARD_TOP, cw, CARD_H, RGBColor(r, g, b))
        # Colored header band
        br, bg_c, bb = int(band_fc[1:3], 16), int(band_fc[3:5], 16), int(band_fc[5:7], 16)
        _rect(slide, lx, CARD_TOP, cw, BAND_H, RGBColor(br, bg_c, bb))
        # Header label on band
        _txt(slide, lx + 0.16, CARD_TOP + 0.11, cw - 0.26, BAND_H - 0.12,
             header, 15, bold=True, color=band_tc)
        # Body bullets
        _bullets(slide, lx + 0.20, CARD_TOP + BAND_H + 0.25,
                 cw - 0.35, CARD_H - BAND_H - 0.35,
                 bullets, size=14, color=CHARCOAL)

    _note(slide, (
        "Kikoff spends roughly $100 million annually across more than ten marketing channels. "
        "The challenge is that their current approach relies heavily on last-touch and platform-reported "
        "attribution — meaning they know someone converted, but not which combination of marketing activity "
        "actually caused it."
    ))


def slide_05(prs, diagram_path):
    slide = _content_slide(prs, 5,
                           "Last-Touch Attribution Gives Credit to the Wrong Channels")
    _bullets(slide, 0.50, 1.62, 5.15, 5.15, [
        ("Current state:", True),
        ("•  Last-touch gives all credit to the final ad clicked", False),
        ("•  Each platform self-reports ROAS — all overclaim", False),
        ("", False),
        ("Desired state:", True),
        ("•  Causal attribution at channel × platform level", False),
        ("•  Diminishing returns visible per channel", False),
        ("", False),
        ("Constraints:", True),
        ("•  Bayesian required (linear regression ruled out)", False),
        ("•  Incrementality priors time-limited to test windows", False),
        ("•  Dual DV: conversions and LTV modeled separately", False),
    ], size=14)
    _img(slide, diagram_path, 5.78, 1.65, 7.12, 5.05)
    _note(slide, (
        "Kikoff is a fintech company focused on credit-building and financial wellness subscriptions. "
        "The problem they're trying to solve is one that many consumer brands face at scale: as marketing "
        "spend grows, traditional attribution methods break down."
    ))


def slide_06(prs):
    slide = _content_slide(prs, 6,
                           "21 Months, 13 Channels, 4 Incrementality Tests — After 5 Quality Fixes")
    _bullets(slide, 0.50, 1.62, 5.15, 5.15, [
        ("Spend data:", True),
        ("•  Daily · 15 source groups → 13 canonical channels", False),
        ("•  Q3 2024 – Mar 2026 (639 rows · 28 cols · 0 null LTV)", False),
        ("", False),
        ("Outcome data:", True),
        ("•  Daily conversions + model-predicted LTV", False),
        ("•  Mean 2,560 conversions/day · peak 4,354", False),
        ("", False),
        ("Incrementality tests:", True),
        ("•  TikTok · Meta ×2 (incl. Jan 2026 partial) · CTV geo", False),
        ("", False),
        ("5 quality issues resolved:", True),
        ("•  Platform typos · LTV anomaly windows (|z| max = 0.85)", False),
        ("•  Null rows · DSP mismatch · Idle/Others erroneous data removed", False),
    ], size=13)
    _img(slide, os.path.join(EDA_FIGS, "channel_composition_monthly.png"),
         5.78, 1.60, 7.12, 5.12)
    _note(slide, (
        "We have three datasets. First: daily spend broken down by source group and platform — iOS, "
        "Android, and Web. Second: daily outcomes — conversion counts and model-predicted lifetime value. "
        "Third: a CSV of historical incrementality tests covering TikTok, two Meta tests, and a CTV test."
    ))


def slide_07(prs, deliverable_preview_path):
    slide = _content_slide(prs, 7,
                           "The Deliverable: ICAC and ROAS Curves That Anchor the Next Budget Cycle")
    _bullets(slide, 0.50, 1.62, 5.15, 5.15, [
        ("Primary outputs:", True),
        ("•  ICAC + ROAS curves per channel × platform", False),
        ("•  Efficiency at current, +20%, −20% spend levels", False),
        ("", False),
        ("Saturation curves:", True),
        ("•  Where diminishing returns kick in per channel", False),
        ("•  Marginal ROI guardrails for budget reallocation", False),
        ("", False),
        ("Seasonal scenarios:", True),
        ("•  Tax season · holiday · steady-state allocations", False),
        ("•  Budget recommendations with spend-move %", False),
        ("", False),
        ("Output format →", True),
    ], size=13)
    _img(slide, deliverable_preview_path, 5.78, 1.60, 7.12, 5.12)
    _note(slide, (
        "Our primary outputs are all at the channel-by-platform level. First: incremental contribution "
        "estimates — for each channel-platform combination, how much of observed conversions and LTV is "
        "attributable to that spend, net of everything else."
    ))


def slide_08(prs, table_path):
    slide = _content_slide(prs, 8,
                           "PyMC-Marketing Won on Flexibility, Integration, and Fallback")
    _bullets(slide, 0.50, 1.60, 12.20, 0.70, [
        ("5 frameworks evaluated — decision locked and sponsor-approved:", True),
    ], size=14)
    _img(slide, table_path, 0.38, 2.40, 12.42, 4.30)
    _note(slide, (
        "We evaluated five modeling frameworks — Meridian, LightweightMMM, Robyn, a custom PyMC build, "
        "and PyMC-Marketing. We selected PyMC-Marketing for its flexibility, its native support for "
        "incrementality priors, and a clean fallback path to raw PyMC if needed."
    ))


def slide_09(prs, phases_path):
    slide = _content_slide(prs, 9,
                           "Five Phases, Phase 1 Done — Meta Model POC Active Now")
    _bullets(slide, 0.50, 1.62, 12.20, 0.78, [
        ("Phase 1 complete · Phase 2 begins now · Phases 3–5: replication → scenarios → handoff", False),
    ], size=14)
    _img(slide, phases_path, 0.38, 2.48, 12.42, 4.70)
    _note(slide, (
        "We're working in five phases. Phase 1 — business understanding and data audit — is essentially "
        "complete. We've locked the taxonomy, selected the framework, and addressed the data quality issues."
    ))


def slide_10(prs):
    slide = _content_slide(prs, 10,
                           "Meta, TikTok, Google Drive 73% of Spend — Tax Season Peaks Are Real")
    stats = [
        ("$112.87M", "total spend in modeling window"),
        ("73.2%",    "share from Meta · TikTok · Google"),
        ("4,354",    "peak conversions/day  (2025-05-30)"),
        ("0.45",     "spend–conversion correlation\n(lag analysis, full date range)"),
    ]
    y = 1.80
    for val, label in stats:
        _txt(slide, 0.50, y, 2.20, 0.52, val, 22, bold=True, color=NAVY)
        _txt(slide, 2.75, y + 0.08, 2.70, 0.44, label, 12, color=CHARCOAL)
        y += 0.90
    _img(slide, os.path.join(EDA_FIGS, "conversions_and_revenue.png"),
         5.78, 1.60, 7.12, 5.12)
    _note(slide, (
        "On the data review: spend and outcome data are now merged at daily granularity, and the modeling "
        "window is confirmed as Q3 2024 through March 2026. Meta, TikTok, and Google account for roughly "
        "two-thirds of total spend. Seasonality is clear in the data — tax season shows a lift."
    ))


def slide_11(prs):
    slide = _content_slide(prs, 11,
                           "Adstock + Saturation + Bayesian Priors = Causal Attribution")
    _bullets(slide, 0.50, 1.62, 5.15, 5.15, [
        ("Adstock (carry-over):", True),
        ("•  Ad effects decay across days — decay rate estimated per channel", False),
        ("•  Data: peak correlation at lag 0 — decay will be short (see →)", False),
        ("", False),
        ("Hill saturation:", True),
        ("•  Diminishing returns curve shape fit per channel", False),
        ("•  Tells us when more spend stops buying proportional growth", False),
        ("", False),
        ("Incrementality priors:", True),
        ("•  Meta, TikTok, CTV holdout test results injected as priors", False),
        ("•  Applied ONLY within test windows — not retroactively (D004)", False),
    ], size=13)
    _img(slide, os.path.join(EDA_FIGS, "spend_conversion_lag.png"),
         5.78, 1.60, 7.12, 5.12)
    _note(slide, (
        "Two transformations happen before the data enters the model. First, adstock — advertising has "
        "a carry-over effect. An ad someone sees today doesn't just influence today's decision; some "
        "fraction of that influence carries into tomorrow and the day after. Second, saturation — also "
        "called the Hill transformation — captures diminishing returns."
    ))


def slide_12(prs, milestones_path):
    slide = _content_slide(prs, 12,
                           "Three Milestones Separate Us from Kikoff's New Budget Model")
    _img(slide, milestones_path, 0.38, 1.52, 12.42, 5.60)
    _note(slide, (
        "To close: we came in with a broad mandate to measure causal channel contribution. Over the past "
        "several weeks, we refined that into a concrete specification — channel-by-platform, dual outcomes, "
        "Bayesian with time-limited incrementality priors — completed the data foundation, and selected the "
        "framework. All of that is locked and sponsor-approved, ahead of schedule."
    ))


def slide_13(prs):
    slide = _title_slide(
        prs,
        title="Questions?",
        subtitle="Yash Tiwari  ·  Chunjiang Liu  ·  Shivang Bhatt  ·  Yoko He",
    )
    _note(slide, "Thank you. We're happy to take questions.")


# ── Appendix slides ────────────────────────────────────────────────────────────

def slide_14_appendix_index(prs):
    slide = _content_slide(prs, "A", "Appendix: Supporting Detail")
    _bullets(slide, 0.50, 1.68, 5.80, 5.05, [
        ("Roadmap:", True),
        ("A-R1.  Project phases — scope, outputs, status", False),
        ("A-R2.  Phase 2 milestones — what each one means", False),
        ("", False),
        ("Data & EDA:", True),
        ("A1.   Full SOURCE_GROUP → channel mapping with spend shares", False),
        ("A2.   DSP vendor spend distribution (StackAdapt + Appvertiser)", False),
        ("A3.   Channel mix composition over time", False),
        ("A4.   Total daily spend trend", False),
        ("A5.   Spend vs. conversions scatter", False),
        ("A6.   LTV per conversion over time", False),
        ("A7.   Seasonality index", False),
    ], size=13)
    _bullets(slide, 6.80, 1.68, 5.80, 5.05, [
        ("Methodology & Reference:", True),
        ("A8.   Framework comparison matrix — full detail", False),
        ("A9.   Incrementality test summary — iCAC, CI, conviction", False),
        ("A10.  Data quality register — findings and resolutions", False),
        ("A11.  MMM concept primer", False),
        ("", False),
        ("Slide References:", True),
        ("A-R1 → next slide", False),
        ("A-R2 → slide after A-R1", False),
        ("EDA charts → A-EDA-1 through A-EDA-3", False),
    ], size=13, color=NAVY)
    _note(slide, "Reference only.")


def slide_15_appendix_eda1(prs):
    slide = _content_slide(prs, "A-EDA-1",
                           "Channel Spend: Top 10 Channels Drive the Model")
    _img(slide, os.path.join(EDA_FIGS, "top_10_channels_spend.png"),
         0.38, 1.52, 12.55, 5.65)
    _note(slide, "Top 10 source groups by total spend in modeling window. Meta (Facebook), TikTok, and Google together account for ~73% of total.")


def slide_16_appendix_eda2(prs):
    slide = _content_slide(prs, "A-EDA-2",
                           "Spend & Conversion Relationship Holds Across Channels")
    _img(slide, os.path.join(EDA_FIGS, "spend_vs_conversions_scatter.png"),
         0.38, 1.52, 6.10, 5.30)
    _txt(slide, 0.38, 6.86, 6.10, 0.30,
         "r = 0.317  ·  modeling window  ·  same-day  (lag analysis r = 0.45, full date range — see slide 11)",
         9, italic=True, color=CHARCOAL)
    _img(slide, os.path.join(EDA_FIGS, "ltv_per_conversion.png"),
         6.55, 1.52, 6.40, 5.65)
    _note(slide, "Left: spend vs conversions scatter, r=0.317 (modeling window, same-day). Lag analysis r=0.45 (full date range, at lag 0) is shown on slide 11. Right: LTV per conversion over time shows mean reversion after anomaly windows were corrected.")


def slide_17_appendix_eda3(prs):
    slide = _content_slide(prs, "A-EDA-3",
                           "Seasonality Is Real — Tax Season and Mid-Year Show Distinct Patterns")
    _img(slide, os.path.join(EDA_FIGS, "monthly_seasonality.png"),
         0.38, 1.52, 6.10, 5.65)
    _img(slide, os.path.join(EDA_FIGS, "total_daily_spend.png"),
         6.55, 1.52, 6.40, 5.65)
    _note(slide, "Left: monthly seasonality index — clear peaks in tax season (Jan-Apr). Right: total daily spend trend showing growth in modeling window.")


def slide_18_appendix_roadmap_phases(prs, phases_table_path):
    slide = _content_slide(prs, "A-R1",
                           "Four Phases: Data Foundation → Model → Replicate → Deliver")
    _txt(slide, 0.50, 1.55, 12.30, 0.50,
         "Green = complete  ·  Yellow = in progress  ·  White = not started", 12,
         color=CHARCOAL, italic=True)
    _img(slide, phases_table_path, 0.38, 2.10, 12.55, 5.00)
    _note(slide, "Full project arc: 4 phases from business understanding through final handoff. "
                 "Phase 1 complete, Phase 2 in progress now.")


def slide_19_appendix_roadmap_milestones(prs, milestones_table_path):
    slide = _content_slide(prs, "A-R2",
                           "Phase 2 Has Six Milestones — M1 and M2 Complete, M3 Active")
    _txt(slide, 0.50, 1.55, 12.30, 0.50,
         "Green = complete  ·  Yellow = in progress  ·  White = queued", 12,
         color=CHARCOAL, italic=True)
    _img(slide, milestones_table_path, 0.38, 2.10, 12.55, 5.00)
    _note(slide, "Phase 2 milestone breakdown. M1 (data cleaning) and M2 (time-limited prior mechanism) "
                 "are complete. M3 (Meta POC model) is the current active milestone.")


# ══════════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("Fast Build — Kikoff MMM Midterm Deck (Run B)")
    print(f"Output: {OUTPUT}")
    print("=" * 60)

    assert _HERE == os.path.join(ROOT, "deliverables", "midterm_presentation", "fast_build"), \
        "OUTPUT is not inside fast_build/ — aborting."

    print("\n[1/2] Generating assets...")
    diagram_path             = _gen_problem_diagram()
    table_path               = _gen_framework_table()
    phases_path              = _gen_phases()
    milestones_path          = _gen_milestones()
    deliverable_preview_path = _gen_deliverable_preview()
    appendix_phases_path     = _gen_appendix_phases()
    appendix_milestones_path = _gen_appendix_milestones()

    print("\n[2/2] Building slides...")
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs, diagram_path)
    slide_06(prs)
    slide_07(prs, deliverable_preview_path)
    slide_08(prs, table_path)
    slide_09(prs, phases_path)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs, milestones_path)
    slide_13(prs)
    slide_14_appendix_index(prs)
    slide_15_appendix_eda1(prs)
    slide_16_appendix_eda2(prs)
    slide_17_appendix_eda3(prs)
    slide_18_appendix_roadmap_phases(prs, appendix_phases_path)
    slide_19_appendix_roadmap_milestones(prs, appendix_milestones_path)

    prs.save(OUTPUT)
    print(f"\n  Saved: {OUTPUT}  ({len(prs.slides)} slides)")

    from pptx import Presentation as _P
    check = _P(OUTPUT)
    assert len(check.slides) == 19, f"Expected 19 slides, got {len(check.slides)}"
    print("  Verification: 19 slides confirmed")


if __name__ == "__main__":
    main()
