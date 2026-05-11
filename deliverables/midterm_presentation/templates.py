"""
Reusable slide template functions for the Kikoff MMM midterm deck.
All geometry and color constants come from design_system.md.
"""

import os
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Colors ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x18, 0x2B, 0x48)
GOLD    = RGBColor(0xFF, 0xCC, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x33, 0x33, 0x33)

# ── Asset paths ──────────────────────────────────────────────────────────────
_BASE   = os.path.dirname(os.path.abspath(__file__))
ASSETS  = os.path.join(_BASE, 'assets')
_CONTENT_TMPL = os.path.join(ASSETS, 'content.png')
_TITLE_TMPL   = os.path.join(ASSETS, 'main_background.png')

# ── Slide dimensions ─────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Layout indices (Blank = 6 in standard pptx layouts) ─────────────────────
_BLANK_LAYOUT_IDX = 6


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT_IDX])


def _add_rect(slide, left, top, width, height, fill_rgb, line=False):
    """Add a solid-filled rectangle with no visible border."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if not line:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, text, left, top, width, height,
                 font_size=16, bold=False, color=None,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    """Add a single-paragraph text box."""
    color = color or CHARCOAL
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return txb


def _set_page_number(slide, page_num: int):
    """Cover the baked-in '2' in content.png and write the correct number."""
    # Navy rectangle to erase the baked-in "2"
    _add_rect(slide, left=0.15, top=7.15, width=0.30, height=0.32,
              fill_rgb=NAVY)
    # Correct page number in gold
    _add_textbox(slide, str(page_num),
                 left=0.17, top=7.16, width=0.28, height=0.28,
                 font_size=9, bold=False, color=GOLD, align=PP_ALIGN.LEFT)


# ── Public template functions ────────────────────────────────────────────────

def add_content_slide(prs, title: str, page_num: int):
    """
    White content slide with navy/gold frame (baked into content.png).
    Returns the slide object for further content addition.
    """
    slide = _blank_slide(prs)

    # Background template (wordmark already inside)
    slide.shapes.add_picture(_CONTENT_TMPL,
                              Inches(0), Inches(0), SLIDE_W, SLIDE_H)

    # Fix page number
    _set_page_number(slide, page_num)

    # Slide title
    _add_textbox(slide, title,
                 left=0.50, top=0.87, width=9.50, height=0.65,
                 font_size=22, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    return slide


def add_title_section_slide(prs, title: str = '', subtitle: str = ''):
    """
    Full navy title/section slide (wordmark baked into main_background.png).
    Returns the slide object for further content addition.
    """
    slide = _blank_slide(prs)

    # Background template
    slide.shapes.add_picture(_TITLE_TMPL,
                              Inches(0), Inches(0), SLIDE_W, SLIDE_H)

    if title:
        _add_textbox(slide, title,
                     left=1.00, top=2.00, width=9.50, height=1.80,
                     font_size=36, bold=True, color=WHITE,
                     align=PP_ALIGN.LEFT)

    if subtitle:
        _add_textbox(slide, subtitle,
                     left=1.00, top=3.90, width=9.50, height=0.70,
                     font_size=20, bold=False, color=GOLD,
                     align=PP_ALIGN.LEFT)

    return slide


def add_bullet_body(slide, bullets: list[str],
                    left=0.50, top=1.65, width=12.20, height=5.00,
                    font_size=16, color=None, bold_first=False):
    """
    Add a bulleted list to an existing content slide.
    bullets: list of strings. Prefix with '\t' for sub-bullet indentation.
    """
    color = color or CHARCOAL
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True

    for i, text in enumerate(bullets):
        indent = text.startswith('\t')
        text = text.lstrip('\t')
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if indent:
            p.level = 1
        run = p.add_run()
        run.text = ('  • ' if indent else '• ') + text
        run.font.size = Pt(font_size - 2 if indent else font_size)
        run.font.bold = (bold_first and i == 0)
        run.font.color.rgb = color
        run.font.name = 'Calibri'

    return txb


def add_two_column_text(slide, left_header, left_items,
                        right_header, right_items,
                        top=1.65, height=5.00, font_size=15):
    """
    Two equal text columns on a content slide.
    """
    col_w = 5.80
    gap = 0.60
    for header, items, left in [
        (left_header, left_items, 0.50),
        (right_header, right_items, 0.50 + col_w + gap)
    ]:
        txb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(col_w), Inches(height)
        )
        tf = txb.text_frame
        tf.word_wrap = True
        # Header paragraph
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = header
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = NAVY
        run.font.name = 'Calibri'
        # Item paragraphs
        for item in items:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = '• ' + item
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.color.rgb = CHARCOAL
            run.font.name = 'Calibri'
