"""
Executive Summary — Kikoff MMM (6 slides)
Audience: UCSD faculty (capstone delivery). Voice: "our team", no internal references.
Reuses the midterm deck's design language (UCSD backgrounds, navy/gold palette,
Calibri, 13.33x7.5 widescreen) and the existing matplotlib asset patterns.
Latest model outputs (P2_04_full_channel, 2026-05-11) carry the "what we built"
and "findings" slides.

Run from project root:
    my-notebook-project/.venv/bin/python deliverables/executive_summary/build_exec_summary.py
"""
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────────────────
_HERE        = os.path.dirname(os.path.abspath(__file__))
ROOT         = os.path.abspath(os.path.join(_HERE, "..", ".."))
MIDTERM_DIR  = os.path.join(ROOT, "deliverables", "midterm_presentation")
ASSETS_SRC   = os.path.join(MIDTERM_DIR, "assets")
MIDTERM_FAST = os.path.join(MIDTERM_DIR, "fast_build", "assets")
EXEC_ASSETS  = os.path.join(_HERE, "assets")
EDA_FIGS     = os.path.join(ROOT, "outputs", "P1_02_eda", "figures")
P204_FIGS    = os.path.join(ROOT, "outputs", "P2_04_full_channel", "figures")
OUTPUT       = os.path.join(_HERE, "kikoff_mmm_executive_summary.pptx")

BG_CONTENT   = os.path.join(ASSETS_SRC, "content.png")
BG_TITLE     = os.path.join(ASSETS_SRC, "main_background.png")

# Reused midterm assets (still valid)
PROBLEM_DIAGRAM = os.path.join(MIDTERM_FAST, "slide_05_problem_diagram.png")
FRAMEWORK_TABLE = os.path.join(MIDTERM_FAST, "slide_08_framework_table.png")
PHASES_DIAGRAM  = os.path.join(MIDTERM_FAST, "slide_09_phases.png")

os.makedirs(EXEC_ASSETS, exist_ok=True)

# ── Brand colors (match midterm exactly) ──────────────────────────────────────
NAVY     = RGBColor(0x18, 0x2B, 0x48)
GOLD     = RGBColor(0xFF, 0xCC, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x33, 0x33, 0x33)
LGRAY    = RGBColor(0xF0, 0xF0, 0xF0)
RED_RGB  = RGBColor(0xCC, 0x00, 0x00)

NAVY_HEX     = "#182B48"
GOLD_HEX     = "#FFCC00"
CHARCOAL_HEX = "#333333"
LGRAY_HEX    = "#F0F0F0"


# ══════════════════════════════════════════════════════════════════════════════
# Slide / text helpers (same API as midterm build_deck.py)
# ══════════════════════════════════════════════════════════════════════════════

def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _txt(slide, left, top, w, h, text, size, bold=False, color=None,
         align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tb.fill.background()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name   = "Calibri"
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color if color else CHARCOAL
    return tb


def _bullets(slide, left, top, w, h, items, size=15, color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        text, bold = item if isinstance(item, tuple) else (item, False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = text
        r.font.name  = "Calibri"
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = color if color else CHARCOAL
        first = False
    return tb


def _note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _rect(slide, left, top, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    bc = border_color if border_color else fill_color
    shape.line.color.rgb = bc
    return shape


def _img(slide, path, left, top, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))
    else:
        _txt(slide, left, top, w, h, f"[image not found: {os.path.basename(path)}]",
             11, color=RED_RGB)


def _link(slide, left, top, w, h, text, url, size=11, bold=False, color=None,
          align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tb.fill.background()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color if color else NAVY
    r.hyperlink.address = url
    return tb


DASHBOARD_URL = "https://capstonekikoff-cni9u6vowxfcgwy8ddmbtw.streamlit.app/"
DASHBOARD_DISPLAY = "capstonekikoff-cni9u6vowxfcgwy8ddmbtw.streamlit.app"


def _title_slide(prs, title, subtitle=None, detail=None):
    slide = _blank_slide(prs)
    _img(slide, BG_TITLE, 0, 0, 13.333, 7.5)
    _txt(slide, 1.00, 2.00, 9.50, 2.00, title, 34, bold=True, color=WHITE)
    if subtitle:
        _txt(slide, 1.00, 4.10, 9.50, 0.70, subtitle, 20, color=GOLD)
    if detail:
        _txt(slide, 1.00, 4.90, 9.50, 1.80, detail, 16, color=WHITE)
    return slide


def _content_slide(prs, page_num, title):
    slide = _blank_slide(prs)
    _img(slide, BG_CONTENT, 0, 0, 13.333, 7.5)
    _rect(slide, 0.13, 7.12, 0.34, 0.34, NAVY)
    _txt(slide, 0.14, 7.13, 0.32, 0.30, str(page_num), 9, color=GOLD)
    _txt(slide, 0.50, 0.87, 12.30, 0.65, title, 21, bold=True, color=NAVY)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# New assets (only what midterm doesn't already cover)
# ══════════════════════════════════════════════════════════════════════════════

def _gen_results_strip():
    """Three-row results strip — Validated / Calibrating / Next."""
    path = os.path.join(EXEC_ASSETS, "exec_results_strip.png")
    fig, ax = plt.subplots(figsize=(12, 5.4), facecolor="white")
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 5.4)
    ax.axis("off")

    rows = [
        (
            "[VALIDATED]   Outputs accepted by client",
            "Saturation curve shape approved   *   8-column decisioning template approved\n"
            "Streamlit dashboard live with 94% HDI confidence bands on iCAC, iROAS, and LTV",
            "#E8F5E9", "#2E7D32",
        ),
        (
            "[CALIBRATING]   Lift-test priors narrowing toward truth",
            "Meta & TikTok priors moving to Conversion Lift Study standards   *   CTV priors remain wider\n"
            "Test-window comparison replaces full-aggregate benchmark check",
            "#FFF9C4", NAVY_HEX,
        ),
        (
            "[NEXT]   Full-channel replication & seasonal scenarios",
            "Validated output set rolling out across remaining channels\n"
            "Seasonal allocations (tax / holiday / steady-state)   *   Handoff package by June 6",
            "#E3F2FD", "#1565C0",
        ),
    ]

    row_h = 1.55
    gap = 0.15
    y_top_start = 5.25
    for idx, (header, body, fc, tc) in enumerate(rows):
        y_top = y_top_start - idx * (row_h + gap)
        r, g, b = int(fc[1:3], 16), int(fc[3:5], 16), int(fc[5:7], 16)
        rect = FancyBboxPatch((0.2, y_top - row_h + 0.10), 11.6, row_h - 0.10,
                              boxstyle="round,pad=0.10",
                              facecolor=(r / 255, g / 255, b / 255),
                              edgecolor=LGRAY_HEX, linewidth=1.2)
        ax.add_patch(rect)
        ax.text(0.55, y_top - 0.35, header, fontsize=14, fontweight="bold",
                color=tc, va="center")
        ax.text(0.55, y_top - 1.00, body, fontsize=11, color=CHARCOAL_HEX,
                va="center", linespacing=1.5)

    fig.tight_layout(pad=0.3)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  Generated: {path}")
    return path


def _gen_build_stats():
    """Compact 'what we built' stats panel — paired with a dashboard chart."""
    path = os.path.join(EXEC_ASSETS, "exec_build_stats.png")

    fig, ax = plt.subplots(figsize=(5.8, 6.0), facecolor="white")
    ax.set_xlim(0, 5.8)
    ax.set_ylim(0, 6.0)
    ax.axis("off")

    stats = [
        ("19", "channels modeled",
         "All paid channels in one model — no subsets"),
        ("93", "weekly observations",
         "Jul 2024 – Mar 2026, weekly (W-MON) rollup"),
        ("2", "models fit in parallel",
         "Model 1: Conversions  ·  Model 2: 3-yr LTV"),
        ("7", "lift tests integrated",
         "Meta x3 (May 25), TikTok x3 (Aug 25), CTV (Oct 25)"),
        ("1.01", "R-hat (Model 2)",
         "Posterior chains within convergence gate (< 1.05)"),
    ]

    y = 5.55
    for big, label, sub in stats:
        ax.text(0.20, y, big, fontsize=30, fontweight="bold", color=NAVY_HEX,
                va="top", ha="left")
        ax.text(1.80, y - 0.08, label, fontsize=11.5, color=CHARCOAL_HEX,
                fontweight="bold", va="top")
        ax.text(1.80, y - 0.55, sub, fontsize=9.5, color=CHARCOAL_HEX,
                va="top", style="italic")
        y -= 1.08

    fig.tight_layout(pad=0.2)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  Generated: {path}")
    return path


def _gen_remaining_work():
    """Path-to-final visual — 3 horizontal cards over the timeline."""
    path = os.path.join(EXEC_ASSETS, "exec_remaining_work.png")

    fig, ax = plt.subplots(figsize=(12, 3.4), facecolor="white")
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 3.4)
    ax.axis("off")

    items = [
        ("Calibration", "Tighter lift-test priors\n(Conversion Lift Study sigma)\nre-fit + windowed validation",
         "#FFF9C4", NAVY_HEX),
        ("Replication", "Validated 6-chart output set\nrolled out for all 19\nchannel x platform combos",
         "#E3F2FD", "#1565C0"),
        ("Scenarios + Handoff", "Tax / holiday / steady-state\nbudget allocations  +  full\nhandoff package to Kikoff",
         "#E8F5E9", "#2E7D32"),
    ]

    cw, gap = 3.7, 0.30
    x = 0.30
    for header, body, fc, tc in items:
        r, g, b = int(fc[1:3], 16), int(fc[3:5], 16), int(fc[5:7], 16)
        rect = FancyBboxPatch((x, 0.55), cw, 2.50,
                              boxstyle="round,pad=0.10",
                              facecolor=(r / 255, g / 255, b / 255),
                              edgecolor=LGRAY_HEX, linewidth=1.2)
        ax.add_patch(rect)
        ax.text(x + 0.25, 2.75, header, fontsize=13, fontweight="bold",
                color=tc, va="center")
        ax.text(x + 0.25, 1.55, body, fontsize=10.5, color=CHARCOAL_HEX,
                va="center", linespacing=1.5)
        x += cw + gap

    ax.text(6.0, 0.20, "Final presentation: June 6, 2026",
            ha="center", fontsize=11, color=NAVY_HEX, fontweight="bold",
            style="italic")

    fig.tight_layout(pad=0.3)
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  Generated: {path}")
    return path


# ══════════════════════════════════════════════════════════════════════════════
# Slide builders
# ══════════════════════════════════════════════════════════════════════════════

def slide_01(prs):
    slide = _title_slide(
        prs,
        title="Measuring What Actually Drives\nKikoff's Growth",
        subtitle="UC San Diego Rady MSBA — Capstone 2026  ·  Executive Summary",
        detail="Yash Tiwari  ·  Chunjiang Liu  ·  Shivang Bhatt  ·  Yoko He\nMay 2026",
    )
    _note(slide, (
        "Executive summary of the Kikoff Marketing Mix Model capstone — built by the UC San Diego "
        "Rady MSBA team. Six slides covering problem, approach, what we built, current findings, "
        "and path to final delivery."
    ))


def slide_02(prs):
    slide = _content_slide(prs, 2,
                           "Platform Attribution Costs Kikoff Precision on $112M in Annual Spend")

    CARD_TOP = 1.62
    CARD_H   = 5.30
    BAND_H   = 0.56
    cw, gap  = 3.92, 0.24

    cards = [
        (
            "The Problem",
            [
                "$112.87M in modeling-window spend",
                "Last-touch credits the final click, not the cause",
                "Each platform self-reports ROAS — all overclaim",
            ],
            "#E3F2FD", "#1565C0", WHITE,
        ),
        (
            "What Breaks at Scale",
            [
                "Pulling Meta spend repeatedly hurts growth within a week",
                "Yet platform dashboards show flat attributed conversions",
                "Budget decisions rest on incomplete attribution signals",
            ],
            "#FFF9C4", "#E65100", WHITE,
        ),
        (
            "What We're Solving",
            [
                "Causal attribution at channel x platform level",
                "Diminishing-return curves per channel for budget reallocation",
                "Defensible model anchoring the next planning cycle",
            ],
            "#E8F5E9", "#2E7D32", WHITE,
        ),
    ]

    for i, (header, bullets, card_fc, band_fc, band_tc) in enumerate(cards):
        lx = 0.38 + i * (cw + gap)
        r, g, b = int(card_fc[1:3], 16), int(card_fc[3:5], 16), int(card_fc[5:7], 16)
        _rect(slide, lx, CARD_TOP, cw, CARD_H, RGBColor(r, g, b))
        br, bg_c, bb = int(band_fc[1:3], 16), int(band_fc[3:5], 16), int(band_fc[5:7], 16)
        _rect(slide, lx, CARD_TOP, cw, BAND_H, RGBColor(br, bg_c, bb))
        _txt(slide, lx + 0.16, CARD_TOP + 0.11, cw - 0.26, BAND_H - 0.12,
             header, 15, bold=True, color=band_tc)
        _bullets(slide, lx + 0.20, CARD_TOP + BAND_H + 0.25,
                 cw - 0.35, CARD_H - BAND_H - 0.35,
                 bullets, size=14, color=CHARCOAL)

    _note(slide, (
        "Kikoff spends roughly $112M annually across more than ten marketing channels. "
        "Current attribution relies on last-touch and platform-reported numbers — both of which "
        "overclaim. The operational signal that gave the team conviction this matters: every time "
        "Meta spend is pulled back, growth slows within a week, yet platform dashboards don't show "
        "that swing. Closing that gap is what the model exists to do."
    ))


def slide_03(prs):
    slide = _content_slide(prs, 3,
                           "Full-Channel Bayesian MMM with Time-Limited Lift-Test Priors")
    _bullets(slide, 0.50, 1.62, 5.30, 5.30, [
        ("What we built:", True),
        ("•  Bayesian MMM — not linear regression", False),
        ("•  All 19 paid channels in one model (no subsets)", False),
        ("•  Two parallel models: Conversions + 3-year LTV", False),
        ("", False),
        ("Why these choices:", True),
        ("•  Bayesian gives full posterior uncertainty per channel", False),
        ("•  Full-channel avoids omitted-variable bias", False),
        ("•  Dual DV separates volume from value", False),
        ("", False),
        ("Lift-test priors:", True),
        ("•  Meta, TikTok, CTV holdout tests injected as priors", False),
        ("•  Applied ONLY within test windows — not retroactively", False),
        ("•  Sigma calibrated per Conversion Lift Study standards", False),
    ], size=13)
    _img(slide, FRAMEWORK_TABLE, 5.95, 1.95, 6.95, 4.20)
    _txt(slide, 5.95, 6.30, 6.95, 0.40,
         "5 frameworks evaluated — PyMC-Marketing selected for flexibility, native lift integration, and dual-DV support",
         10, italic=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
    _note(slide, (
        "Three methodological commitments anchor the model. First, Bayesian — we need full posterior "
        "uncertainty per channel, not a point estimate, because most channels will end up at medium "
        "or low trust. Second, full-channel scope — modeling a subset would absorb other channels' "
        "effects into included channels' coefficients. Third, dual DV — conversions and 3-year LTV "
        "carry different decisioning signals, so we fit both in parallel. PyMC-Marketing won the "
        "framework evaluation on flexibility, native lift integration, and a clean fallback path."
    ))


def slide_04(prs, build_stats_path):
    slide = _content_slide(prs, 4,
                           "Two Models Fit on All 19 Channels — Live Dashboard with Confidence Bands")
    _img(slide, build_stats_path, 0.38, 1.60, 5.70, 5.20)

    # Hero chart: latest saturation curve from May 11 model output
    _img(slide, os.path.join(P204_FIGS, "meta_web_saturation_curve.png"),
         6.30, 1.65, 6.70, 4.15)

    # Caption for the chart
    _txt(slide, 6.30, 5.88, 6.70, 0.40,
         "Meta Web saturation curve — 94% HDI band, vertical line at median weekly spend",
         10, italic=True, color=CHARCOAL, align=PP_ALIGN.CENTER)

    # Dashboard CTA — full-width band at the bottom
    _rect(slide, 0.38, 6.55, 12.55, 0.60, NAVY)
    _txt(slide, 0.50, 6.62, 4.40, 0.45,
         "Live dashboard ▸",
         13, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    _link(slide, 2.30, 6.65, 10.50, 0.42,
          DASHBOARD_DISPLAY,
          DASHBOARD_URL,
          size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    _note(slide, (
        "Phase 2 execution. The model is fit and the dashboard is live. Both models — Conversions "
        "and 3-year LTV — run on all 19 channels at weekly granularity across 93 weeks. Seven lift "
        "tests are integrated as windowed priors, applying only within each test's own date range. "
        "The saturation curve shown on the right is the kind of output the client uses for budget "
        "decisioning — spend on the X-axis, incremental CAC on the Y-axis, posterior band shaded, "
        "and a vertical line marking current weekly spend. This is one of eight chart types live in "
        "the Streamlit dashboard, hosted at the URL shown — updates automatically as the model is "
        "recalibrated and new outputs are pushed."
    ))


def slide_05(prs, results_strip_path):
    slide = _content_slide(prs, 5,
                           "Saturation Shape & Decisioning Layer Approved; Calibration Tightening Next")
    # Results strip is the hero — full working area, no crammed supporting charts
    _img(slide, results_strip_path, 0.50, 1.75, 12.33, 5.10)

    # Subtle dashboard pointer at the very bottom
    _txt(slide, 0.50, 6.95, 6.50, 0.35,
         "Explore every chart referenced here on the live dashboard:",
         11, italic=True, color=CHARCOAL, align=PP_ALIGN.RIGHT)
    _link(slide, 7.05, 6.95, 5.80, 0.35,
          DASHBOARD_DISPLAY,
          DASHBOARD_URL,
          size=11, bold=True, color=NAVY, italic=True, align=PP_ALIGN.LEFT)

    _note(slide, (
        "Three-row read on current state. First, what's validated: the client signed off on the "
        "saturation curve shape and on the 8-column decisioning template, both reviewed live in the "
        "dashboard walkthrough. Second, what's calibrating: lift-test sigmas on Meta and TikTok are "
        "moving from an initial heuristic to Conversion Lift Study standards, which will tighten "
        "posterior contributions for those channels; CTV's sigma stays wider because its CSV "
        "provides a true confidence interval. Benchmark comparisons are also shifting from "
        "full-history aggregates to test-window-specific reads, since that's where the lift test's "
        "truth applies. Third, what's next: roll the validated output set out across all channels "
        "and start the seasonal-scenario work. Every chart referenced is live on the dashboard "
        "URL at the bottom."
    ))


def slide_06(prs, remaining_work_path):
    slide = _content_slide(prs, 6,
                           "Three Workstreams Carry Us to Final Delivery on June 6")
    _img(slide, PHASES_DIAGRAM, 0.38, 1.55, 12.55, 2.40)
    _img(slide, remaining_work_path, 0.38, 4.00, 12.55, 3.10)

    _note(slide, (
        "Closing read. Phase 1 — data foundation and framework selection — is complete. Phase 2 — "
        "model build, dashboard, decisioning template — is in flight and validated in form. Three "
        "workstreams remain between now and June 6: calibration (tightening lift-test priors and "
        "re-fitting), replication (rolling the validated output set across all 19 channel x platform "
        "combinations), and scenario planning plus handoff. The model architecture is locked, the "
        "deliverable shape is locked, and the path to the final presentation is concrete."
    ))


# ══════════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("Executive Summary — Kikoff MMM (6 slides)")
    print(f"Output: {OUTPUT}")
    print("=" * 60)

    assert _HERE == os.path.join(ROOT, "deliverables", "executive_summary"), \
        "OUTPUT is not inside executive_summary/ — aborting."

    # Verify reused midterm assets exist
    for required in [BG_CONTENT, BG_TITLE, FRAMEWORK_TABLE, PHASES_DIAGRAM]:
        assert os.path.exists(required), f"Missing reused asset: {required}"

    # Verify latest model outputs exist
    for required in [
        os.path.join(P204_FIGS, "meta_web_saturation_curve.png"),
        os.path.join(P204_FIGS, "meta_web_icac_over_time.png"),
        os.path.join(P204_FIGS, "meta_web_iroas_saturation_curve.png"),
    ]:
        assert os.path.exists(required), f"Missing model output: {required}"

    print("\n[1/2] Generating new assets...")
    results_strip_path  = _gen_results_strip()
    build_stats_path    = _gen_build_stats()
    remaining_work_path = _gen_remaining_work()

    print("\n[2/2] Building slides...")
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs, build_stats_path)
    slide_05(prs, results_strip_path)
    slide_06(prs, remaining_work_path)

    prs.save(OUTPUT)
    print(f"\n  Saved: {OUTPUT}  ({len(prs.slides)} slides)")

    from pptx import Presentation as _P
    check = _P(OUTPUT)
    assert len(check.slides) == 6, f"Expected 6 slides, got {len(check.slides)}"
    print("  Verification: 6 slides confirmed")


if __name__ == "__main__":
    main()
